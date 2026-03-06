"""
plugins/io_actions.py
Azioni per creare file/cartelle, scrivere contenuti, ricerca file, impostazioni Windows.
fix_file ora usa Gemma locale tramite llama-cpp invece di OpenAI.
"""

from pathlib import Path
import os
import sys
import time
import shutil
import subprocess
import ctypes
import ctypes.wintypes
import datetime
import re
import getpass
from typing import Dict, Any, List


# ===================================================
# UTILITY
# ===================================================



def get_real_desktop() -> Path:
    try:
        buf = ctypes.create_unicode_buffer(ctypes.wintypes.MAX_PATH)
        ctypes.windll.shell32.SHGetFolderPathW(None, 0x10, None, 0, buf)
        p = Path(buf.value)
        if p.exists():
            return p
    except Exception:
        pass
    p1 = Path.home() / "OneDrive" / "Desktop"
    p2 = Path.home() / "Desktop"
    return p1 if p1.exists() else p2


def normalize_path(path: str, workspace: Path = None) -> Path:
    """
    Risolve path semplici che l'AI può passare direttamente:
      desktop/file.txt      → vero Desktop (OneDrive o normale)
      downloads/file.txt    → cartella Downloads
      documents/file.txt    → cartella Documenti
      pictures/file.txt     → cartella Immagini
      music/file.txt        → cartella Musica
      videos/file.txt       → cartella Video
      home/file.txt         → cartella utente
      ~ / percorsi assoluti → risolti normalmente
    """
    if path is None:
        raise ValueError("path is None")

    try:
        username = os.getlogin()
    except Exception:
        username = os.getenv("USERNAME") or ""

    path = path.replace("<USERNAME>", username).replace("%USERNAME%", username)
    path = os.path.expandvars(path)
    path = os.path.expanduser(path)
    p = path.replace("\\", "/").strip()
    lower = p.lower()

    # Alias cartelle comuni — l'AI può usare questi nomi direttamente
    FOLDER_ALIASES = {
        "desktop":   get_real_desktop(),
        "downloads": Path.home() / "Downloads",
        "download":  Path.home() / "Downloads",
        "documents": Path.home() / "Documents",
        "documenti": Path.home() / "Documents",
        "pictures":  Path.home() / "Pictures",
        "immagini":  Path.home() / "Pictures",
        "music":     Path.home() / "Music",
        "musica":    Path.home() / "Music",
        "videos":    Path.home() / "Videos",
        "video":     Path.home() / "Videos",
        "home":      Path.home(),
    }

    # Controlla se il path inizia con un alias (es. "desktop/ciao.txt")
    parts = p.split("/")
    first = parts[0].lower()
    if first in FOLDER_ALIASES:
        base = FOLDER_ALIASES[first]
        tail = "/".join(parts[1:])
        return (base / tail).resolve() if tail else base.resolve()

    # Controlla se contiene "desktop" nel mezzo (es. percorsi parziali Windows)
    if "desktop" in lower:
        desktop = get_real_desktop()
        idx = next((i for i, part in enumerate(parts) if part.lower() == "desktop"), None)
        if idx is not None:
            tail = "/".join(parts[idx + 1:])
            return (desktop / tail).resolve() if tail else desktop.resolve()

    cand = Path(p)
    if cand.is_absolute():
        return cand.resolve()
    if workspace:
        return (workspace / cand).resolve()
    return cand.resolve()


def _is_protected_path(p: Path) -> bool:
    try:
        roots = [Path(os.environ.get("SystemRoot", "C:\\Windows")).resolve()]
        if sys.platform == "win32":
            roots += [
                Path(r"C:\Windows\System32").resolve(),
                Path(r"C:\Program Files").resolve(),
                Path(r"C:\Program Files (x86)").resolve(),
            ]
        for r in roots:
            try:
                if r in p.parents or p == r:
                    return True
            except Exception:
                continue
    except Exception:
        pass
    return False


_blocked_ext = {".dll", ".sys", ".exe", ".msi", ".drv"}


# ===================================================
# BASIC IO
# ===================================================

def create_folder(args: Dict[str, Any], workspace: Path, safe_resolve) -> Dict[str, Any]:
    path = args.get("path")
    if not path:
        return {"error": "missing_path"}
    p = normalize_path(path, workspace)
    if _is_protected_path(p):
        return {"error": "protected_path", "path": str(p)}
    p.mkdir(parents=True, exist_ok=True)
    return {"status": "ok", "created": str(p)}


def create_file(args: Dict[str, Any], workspace: Path, safe_resolve) -> Dict[str, Any]:
    path = args.get("path")
    content = args.get("content", "")
    if not path:
        return {"error": "missing_path"}
    p = normalize_path(path, workspace)
    if _is_protected_path(p):
        return {"error": "protected_path", "path": str(p)}
    if p.suffix.lower() in _blocked_ext:
        return {"error": "blocked_extension", "extension": p.suffix}
    p.parent.mkdir(parents=True, exist_ok=True)
    if p.exists():
        return {"error": "exists", "path": str(p)}
    try:
        p.write_text(content, encoding="utf-8")
    except Exception as e:
        return {"error": "write_failed", "reason": str(e)}
    return {"status": "ok", "created": str(p)}


def write_file(args: Dict[str, Any], workspace: Path, safe_resolve) -> Dict[str, Any]:
    path = args.get("path")
    content = args.get("content", "")
    mode = args.get("mode", "w")
    if not path:
        return {"error": "missing_path"}
    p = normalize_path(path, workspace)
    if _is_protected_path(p):
        return {"error": "protected_path", "path": str(p)}
    if p.suffix.lower() in _blocked_ext:
        return {"error": "blocked_extension", "extension": p.suffix}
    p.parent.mkdir(parents=True, exist_ok=True)
    try:
        with p.open(mode, encoding="utf-8") as fh:
            fh.write(content)
    except Exception as e:
        return {"error": "write_failed", "reason": str(e)}
    return {"status": "ok", "written": str(p)}


# ===================================================
# FIX FILE (ora usa Gemma locale)
# ===================================================

def fix_file(args: Dict[str, Any], workspace: Path, safe_resolve) -> Dict[str, Any]:
    """
    Modifica o corregge un file usando Gemma locale (llama-cpp).
    - args.changes: descrizione delle modifiche da applicare
    - args.content: sostituisci completamente con questo testo
    """
    path_in = args.get("path") or args.get("file_path")
    if not path_in:
        return {"error": "missing_path"}

    p = normalize_path(path_in, workspace)
    if _is_protected_path(p):
        return {"error": "protected_path", "path": str(p)}
    if p.suffix.lower() in _blocked_ext:
        return {"error": "blocked_extension", "extension": p.suffix}
    if not p.exists():
        return {"error": "not_found", "path": str(p)}

    try:
        orig_text = p.read_text(encoding="utf-8")
    except Exception as e:
        return {"error": "not_text", "reason": str(e)}

    # Se viene passato content diretto, sostituisce senza AI
    content = args.get("content")
    if content:
        ts = time.strftime("%Y%m%d_%H%M%S")
        backup = p.with_name(p.name + f".bak_{ts}")
        shutil.copy2(p, backup)
        p.write_text(content, encoding="utf-8")
        return {"status": "ok", "result": "File sostituito", "backup": str(backup)}

    changes = args.get("changes", "")
    if not changes:
        return {"error": "missing_changes", "reason": "Specifica 'changes' o 'content'"}

    # Carica Gemma
    try:
        from llama_cpp import Llama
        gguf_files = list(Path("./models").rglob("*.gguf"))
        if not gguf_files:
            return {"error": "model_not_found", "reason": "Nessun GGUF in ./models"}

        llm = Llama(model_path=str(gguf_files[0]), n_ctx=4096, n_gpu_layers=-1, verbose=False)
    except Exception as e:
        return {"error": "model_load_failed", "reason": str(e)}

    prompt = (
        f"Applica le seguenti modifiche al file {p.name}:\n{changes}\n\n"
        f"Contenuto attuale:\n{orig_text}\n\n"
        "Rispondi SOLO con il nuovo contenuto completo del file, senza spiegazioni."
    )

    try:
        resp = llm.create_chat_completion(
            messages=[
                {"role": "system", "content": "Sei un assistente che modifica file. Rispondi solo con il nuovo contenuto del file."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=2048,
            temperature=0.1,
        )
        new_content = resp["choices"][0]["message"]["content"].strip()
        # Rimuovi eventuali backtick
        new_content = re.sub(r"^```[a-zA-Z]*\n?", "", new_content)
        new_content = re.sub(r"\n?```$", "", new_content).strip()
    except Exception as e:
        return {"error": "model_failed", "reason": str(e)}

    if not new_content:
        return {"error": "empty_response"}

    ts = time.strftime("%Y%m%d_%H%M%S")
    backup = p.with_name(p.name + f".bak_{ts}")
    shutil.copy2(p, backup)
    p.write_text(new_content, encoding="utf-8")
    return {"status": "ok", "result": "File modificato", "backup": str(backup), "path": str(p)}


# ===================================================
# MEDIA
# ===================================================

def _open_with_default(path: Path) -> bool:
    try:
        if sys.platform.startswith("win"):
            os.startfile(str(path))
        elif sys.platform == "darwin":
            subprocess.run(["open", str(path)], check=False)
        else:
            subprocess.run(["xdg-open", str(path)], check=False)
        return True
    except Exception:
        return False


def _guess_open_command_for_ext(ext: str) -> List[List[str]]:
    ext = ext.lower()
    candidates: List[List[str]] = []
    if ext in {".jpg", ".jpeg", ".png", ".bmp", ".gif", ".webp"}:
        candidates += [["mspaint"], ["explorer"], ["xdg-open"], ["open"]]
    elif ext in {".mp4", ".mkv", ".avi", ".mov"}:
        candidates += [["vlc"], ["mpv"], ["explorer"], ["xdg-open"], ["open"]]
    elif ext in {".mp3", ".wav", ".flac", ".aac"}:
        candidates += [["vlc"], ["mpv"], ["xdg-open"], ["open"]]
    elif ext == ".pdf":
        candidates += [["sumatrapdf"], ["acrord32"], ["xdg-open"], ["open"]]
    else:
        candidates += [["explorer"], ["xdg-open"], ["open"]]
    return candidates


def _try_open_with_candidates(path: Path, candidates: List[List[str]]) -> bool:
    for cmd in candidates:
        try:
            subprocess.run(cmd + [str(path)], check=False)
            return True
        except FileNotFoundError:
            continue
    return False


def open_media(args: Dict[str, Any], workspace: Path, safe_resolve) -> Dict[str, Any]:
    path = args.get("path")
    pattern = args.get("pattern")
    recursive = bool(args.get("recursive", True))
    max_open = int(args.get("max_open", 20))
    found: List[Path] = []

    if pattern and not path:
        roots = [Path.home() / "Pictures", Path.home() / "Downloads", get_real_desktop(), Path.home()]
        exts = ["*.jpg", "*.jpeg", "*.png", "*.gif", "*.mp4", "*.mkv", "*.mp3", "*.wav", "*.pdf"]
        for root in roots:
            if not root or not root.exists():
                continue
            for ext in exts:
                for f in (root.rglob(ext) if recursive else root.glob(ext)):
                    if pattern.lower() in f.name.lower() and f.is_file():
                        found.append(f)
    else:
        if not path:
            return {"error": "missing_path_or_pattern"}
        p = normalize_path(path, workspace)
        if p.is_dir():
            exts = ["*.jpg", "*.jpeg", "*.png", "*.gif", "*.mp4", "*.mkv", "*.mp3", "*.wav", "*.pdf"]
            for ext in exts:
                for f in (p.rglob(ext) if recursive else p.glob(ext)):
                    if f.is_file():
                        found.append(f)
        else:
            found.append(p)

    found = list(dict.fromkeys(found))
    if not found:
        return {"status": "ok", "result": "Nessun file trovato"}
    if len(found) > max_open:
        found = found[:max_open]

    opened, failed = [], []
    for f in found:
        if _open_with_default(f):
            opened.append(str(f))
        else:
            ext = f.suffix.lower()
            if _try_open_with_candidates(f, _guess_open_command_for_ext(ext)):
                opened.append(str(f))
            else:
                failed.append(str(f))

    return {"status": "ok", "opened": opened, "failed": failed}


# ===================================================
# SEARCH FILES
# ===================================================

def search_files(args: Dict[str, Any], workspace: Path, safe_resolve) -> Dict[str, Any]:
    query = args.get("query", "")
    name = args.get("name") or args.get("filename")
    ext = args.get("ext") or args.get("extension")
    min_size = args.get("min_size")
    max_size = args.get("max_size")
    limit = int(args.get("limit", 500))
    recursive = bool(args.get("recursive", True))
    root = args.get("path_root") or args.get("path") or str(workspace)

    username = getpass.getuser()
    root = root.replace("<USERNAME>", username)
    if root.startswith("~"):
        root = str(Path.home() / root[2:])

    if query:
        q = query.lower()
        match_name = re.search(r"(?:file|immagine|foto|video|documento)\s+([\w\-. ]+)", q)
        if match_name and not name:
            name = match_name.group(1).strip()
        match_ext = re.search(r"\.(\w+)", q)
        if match_ext and not ext:
            ext = match_ext.group(1).strip()
        if "desktop" in q:
            root = str(get_real_desktop())
        elif "download" in q:
            root = str(Path.home() / "Downloads")
        elif "documenti" in q or "documents" in q:
            root = str(Path.home() / "Documents")
        elif "musica" in q or "music" in q:
            root = str(Path.home() / "Music")
        elif "video" in q:
            root = str(Path.home() / "Videos")

    try:
        root_path = normalize_path(root, workspace)
    except Exception:
        root_path = workspace

    if not root_path.exists():
        return {"error": "root_not_found", "path": str(root_path)}

    exts: List[str] = []
    if ext:
        exts = [ext if ext.startswith(".") else f".{ext}"]

    results = []
    it = root_path.rglob("*") if recursive else root_path.glob("*")
    for p in it:
        if not p.is_file():
            continue
        try:
            if name and name.lower() not in p.name.lower():
                continue
            if exts and p.suffix.lower() not in exts:
                continue
            st = p.stat()
            if min_size and st.st_size < int(min_size):
                continue
            if max_size and st.st_size > int(max_size):
                continue
            results.append({"path": str(p), "size": st.st_size, "modified": time.ctime(st.st_mtime)})
            if len(results) >= limit:
                break
        except Exception:
            continue

    return {"status": "ok", "root": str(root_path), "count": len(results), "results": results}


# ===================================================
# EXPLORE FILES
# ===================================================

def explore_files(args: Dict[str, Any], workspace: Path, safe_resolve) -> Dict[str, Any]:
    mode = args.get("mode", "top_large")
    root = args.get("path_root") or str(Path.home())
    limit = int(args.get("limit", 20))
    hours = int(args.get("hours", 24))
    date_filter = args.get("date_filter", "").lower()
    ext_filter = args.get("ext_filter", "").lower()
    sort_by = args.get("sort", "date").lower()

    root_path = normalize_path(root, workspace)
    if not root_path.exists():
        return {"error": "root_not_found", "path": str(root_path)}

    results = []
    now = datetime.datetime.now()

    EXT_GROUPS = {
        "immagini": [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp"],
        "video": [".mp4", ".mkv", ".avi", ".mov"],
        "audio": [".mp3", ".wav", ".flac", ".aac"],
        "documenti": [".txt", ".docx", ".pdf", ".odt", ".md"],
        "pdf": [".pdf"],
    }

    def match_date(p):
        try:
            mtime = datetime.datetime.fromtimestamp(p.stat().st_mtime)
            if date_filter == "oggi":
                return mtime.date() == now.date()
            elif date_filter == "ieri":
                return mtime.date() == (now - datetime.timedelta(days=1)).date()
            elif date_filter == "settimana":
                return (now - mtime).days <= 7
            elif date_filter == "mese":
                return (now - mtime).days <= 30
            return True
        except Exception:
            return False

    def match_ext(p):
        if not ext_filter:
            return True
        for key, exts in EXT_GROUPS.items():
            if ext_filter in key:
                return p.suffix.lower() in exts
        ext = ext_filter if ext_filter.startswith(".") else f".{ext_filter}"
        return p.suffix.lower() == ext

    if mode == "top_large":
        files = []
        for p in root_path.rglob("*"):
            if p.is_file() and match_date(p) and match_ext(p):
                try:
                    files.append((p.stat().st_size, p))
                except Exception:
                    continue
        files.sort(reverse=True, key=lambda x: x[0])
        for size, p in files[:limit]:
            results.append({"path": str(p), "size": size, "modified": time.ctime(p.stat().st_mtime)})

    elif mode == "recent":
        cutoff = time.time() - (hours * 3600)
        for p in root_path.rglob("*"):
            if p.is_file() and match_date(p) and match_ext(p):
                try:
                    if p.stat().st_mtime >= cutoff:
                        results.append({"path": str(p), "size": p.stat().st_size, "modified": time.ctime(p.stat().st_mtime)})
                        if len(results) >= limit:
                            break
                except Exception:
                    continue

    elif mode == "downloads":
        dl = Path.home() / "Downloads"
        if not dl.exists():
            return {"status": "ok", "results": [], "note": "Downloads non trovata"}
        files = sorted([p for p in dl.iterdir() if p.is_file() and match_ext(p)], key=lambda x: x.stat().st_mtime, reverse=True)
        for p in files[:limit]:
            results.append({"path": str(p), "size": p.stat().st_size, "modified": time.ctime(p.stat().st_mtime)})

    elif mode == "filter":
        for p in root_path.rglob("*"):
            if not p.is_file():
                continue
            try:
                if match_date(p) and match_ext(p):
                    results.append({"path": str(p), "size": p.stat().st_size, "modified": time.ctime(p.stat().st_mtime)})
            except Exception:
                continue
        if sort_by == "size":
            results.sort(key=lambda x: x["size"], reverse=True)
        elif sort_by == "name":
            results.sort(key=lambda x: Path(x["path"]).name.lower())
        else:
            results.sort(key=lambda x: x["modified"], reverse=True)
        results = results[:limit]

    return {"status": "ok", "mode": mode, "root": str(root_path), "count": len(results), "results": results}


# ===================================================
# SYSTEM SETTINGS
# ===================================================

def list_settings_shortcuts(args: Dict[str, Any], workspace: Path, safe_resolve) -> Dict[str, Any]:
    shortcuts = {
        "display": "ms-settings:display", "sound": "ms-settings:sound",
        "notifications": "ms-settings:notifications", "battery": "ms-settings:batterysaver",
        "power": "ms-settings:powersleep", "storage": "ms-settings:storagesense",
        "clipboard": "ms-settings:clipboard", "about": "ms-settings:about",
        "bluetooth": "ms-settings:bluetooth", "printers": "ms-settings:printers",
        "mouse": "ms-settings:mousetouchpad", "typing": "ms-settings:typing",
        "network": "ms-settings:network", "wifi": "ms-settings:network-wifi",
        "ethernet": "ms-settings:network-ethernet", "vpn": "ms-settings:network-vpn",
        "proxy": "ms-settings:network-proxy", "airplane": "ms-settings:network-airplanemode",
        "background": "ms-settings:personalization-background",
        "colors": "ms-settings:personalization-colors", "lockscreen": "ms-settings:lockscreen",
        "themes": "ms-settings:themes", "fonts": "ms-settings:fonts",
        "your_info": "ms-settings:yourinfo", "email": "ms-settings:emailandaccounts",
        "signin": "ms-settings:signinoptions", "date_time": "ms-settings:dateandtime",
        "language": "ms-settings:regionlanguage", "privacy": "ms-settings:privacy",
        "camera": "ms-settings:privacy-webcam", "microphone": "ms-settings:privacy-microphone",
        "location": "ms-settings:privacy-location", "security": "ms-settings:windowsdefender",
        "windows_update": "ms-settings:windowsupdate", "recovery": "ms-settings:recovery",
        "developers": "ms-settings:developers", "remote_desktop": "ms-settings:remotedesktop",
    }
    aliases = {
        "rete": "network", "connessione": "network", "internet": "network",
        "audio": "sound", "suono": "sound", "volume": "sound",
        "schermo": "display", "batteria": "battery",
        "aggiornamento": "windows_update", "update": "windows_update",
        "tastiera": "typing", "notifiche": "notifications",
        "sfondo": "background", "colore": "colors",
        "telecamera": "camera", "microfono": "microphone",
        "orario": "date_time", "ora": "date_time", "lingua": "language",
        "account": "your_info", "sicurezza": "security",
    }
    shortcuts.update({alias: shortcuts[target] for alias, target in aliases.items() if target in shortcuts})
    return {"status": "ok", "shortcuts": shortcuts}


def open_system_setting(args: Dict[str, Any], workspace: Path, safe_resolve) -> Dict[str, Any]:
    key = args.get("key") or args.get("setting")
    if not key:
        return {"error": "missing_key"}
    mapping = list_settings_shortcuts({}, workspace, safe_resolve).get("shortcuts", {})
    uri = mapping.get(key.lower())
    if not uri:
        return {"error": "unknown_key", "key": key}
    try:
        if sys.platform.startswith("win"):
            subprocess.run(["start", "", uri], shell=True, check=False)
            return {"status": "ok", "result": f"Aperto: {key}"}
        else:
            return {"error": "unsupported_platform"}
    except Exception as e:
        return {"error": "open_failed", "reason": str(e)}



# ===================================================
# RENAME FILE
# ===================================================

def rename_file(args: Dict[str, Any], workspace: Path, safe_resolve) -> Dict[str, Any]:
    """Rinomina un file o cartella."""
    path = args.get("path")
    new_name = args.get("new_name")
    if not path:
        return {"error": "missing_path"}
    if not new_name:
        return {"error": "missing_new_name"}
    p = normalize_path(path, workspace)
    if not p.exists():
        return {"error": "not_found", "path": str(p)}
    if _is_protected_path(p):
        return {"error": "protected_path", "path": str(p)}
    new_path = p.parent / new_name
    if new_path.exists():
        return {"error": "already_exists", "path": str(new_path)}
    try:
        p.rename(new_path)
        return {"status": "ok", "renamed": str(new_path)}
    except Exception as e:
        return {"error": "rename_failed", "reason": str(e)}


# ===================================================
# REGISTER
# ===================================================

def register_actions() -> Dict[str, Any]:
    return {
        "create_folder": create_folder,
        "create_file": create_file,
        "write_file": write_file,
        "fix_file": fix_file,
        "rename_file": rename_file,
        "open_media": open_media,
        "open_file": open_media,
        "search_files": search_files,
        "explore_files": explore_files,
        "list_settings_shortcuts": list_settings_shortcuts,
        "open_system_setting": open_system_setting,
        "delete_file": delete_file,
        "delete_folder": delete_folder,
        "copy_file": copy_file,
        "move_file": move_file,
    }

# ===================================================
# .PYTEMP — CESTINO SICURO
# ===================================================

def _get_pytemp_dir() -> Path:
    """Restituisce ~/.pytemp, creandola nascosta se non esiste."""
    pytemp = Path.home() / ".pytemp"
    pytemp.mkdir(exist_ok=True)
    try:
        import ctypes
        ctypes.windll.kernel32.SetFileAttributesW(str(pytemp), 0x02)
    except Exception:
        pass
    return pytemp


def _backup_to_pytemp(path: Path) -> str:
    """Zippa file/cartella in .pytemp con nome [YYYYMMDD_HHMMSS] nomefile.zip"""
    import zipfile
    from datetime import datetime
    pytemp = _get_pytemp_dir()
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    zip_path = pytemp / f"[{timestamp}] {path.name}.zip"
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        if path.is_file():
            zf.write(path, path.name)
        elif path.is_dir():
            for f in path.rglob("*"):
                if f.is_file():
                    zf.write(f, f.relative_to(path.parent))
    return str(zip_path)


def delete_file(args: Dict[str, Any], workspace: Path, safe_resolve) -> Dict[str, Any]:
    """Cancella un file dopo backup in .pytemp."""
    path = args.get("path")
    if not path:
        return {"error": "missing_path"}
    p = normalize_path(path, workspace)
    if not p.exists():
        return {"error": "not_found", "path": str(p)}
    if not p.is_file():
        return {"error": "not_a_file", "path": str(p)}
    if _is_protected_path(p):
        return {"error": "protected_path", "path": str(p)}
    try:
        backup = _backup_to_pytemp(p)
        p.unlink()
        print(f"   '{p.name}' deleted  (backup saved in .pytemp)")
        return {"status": "ok", "deleted": str(p), "backup": backup}
    except Exception as e:
        return {"error": "delete_failed", "reason": str(e)}


def delete_folder(args: Dict[str, Any], workspace: Path, safe_resolve) -> Dict[str, Any]:
    """Cancella una cartella dopo backup in .pytemp."""
    import shutil
    path = args.get("path")
    if not path:
        return {"error": "missing_path"}
    p = normalize_path(path, workspace)
    if not p.exists():
        return {"error": "not_found", "path": str(p)}
    if not p.is_dir():
        return {"error": "not_a_folder", "path": str(p)}
    if _is_protected_path(p):
        return {"error": "protected_path", "path": str(p)}
    try:
        backup = _backup_to_pytemp(p)
        shutil.rmtree(p)
        print(f"   Folder '{p.name}' deleted  (backup saved in .pytemp)")
        return {"status": "ok", "deleted": str(p), "backup": backup}
    except Exception as e:
        return {"error": "delete_failed", "reason": str(e)}


def copy_file(args: Dict[str, Any], workspace: Path, safe_resolve) -> Dict[str, Any]:
    """Copia un file."""
    import shutil
    src = args.get("src")
    dst = args.get("dst")
    if not src or not dst:
        return {"error": "missing_src_or_dst"}
    s = normalize_path(src, workspace)
    d = normalize_path(dst, workspace)
    if not s.exists():
        return {"error": "not_found", "path": str(s)}
    try:
        d.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(str(s), str(d))
        return {"status": "ok", "copied": str(d)}
    except Exception as e:
        return {"error": str(e)}


def move_file(args: Dict[str, Any], workspace: Path, safe_resolve) -> Dict[str, Any]:
    """Sposta un file."""
    import shutil
    src = args.get("src")
    dst = args.get("dst")
    if not src or not dst:
        return {"error": "missing_src_or_dst"}
    s = normalize_path(src, workspace)
    d = normalize_path(dst, workspace)
    if not s.exists():
        return {"error": "not_found", "path": str(s)}
    try:
        d.parent.mkdir(parents=True, exist_ok=True)
        shutil.move(str(s), str(d))
        return {"status": "ok", "moved": str(d)}
    except Exception as e:
        return {"error": str(e)}
# ===================================================
# OFFICE & PDF
# ===================================================

def _check_lib(name: str) -> bool:
    import importlib
    try:
        importlib.import_module(name)
        return True
    except ImportError:
        return False


def create_docx(args: Dict[str, Any], workspace: Path, safe_resolve) -> Dict[str, Any]:
    """
    Crea un file Word (.docx).
    args: {
      "path": "desktop/documento.docx",
      "title": "Titolo del documento",
      "content": [
        {"type": "heading", "text": "Titolo sezione", "level": 1},
        {"type": "paragraph", "text": "Testo normale..."},
        {"type": "bullet", "items": ["voce 1", "voce 2", "voce 3"]},
        {"type": "table", "headers": ["Col1","Col2"], "rows": [["A","B"],["C","D"]]}
      ]
    }
    """
    if not _check_lib("docx"):
        return {"error": "missing_library", "reason": "Installa python-docx: pip install python-docx"}

    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    path = args.get("path", "")
    if not path:
        return {"error": "missing_path"}

    p = normalize_path(path, workspace)
    p.parent.mkdir(parents=True, exist_ok=True)

    doc = Document()
    title = args.get("title", "")
    if title:
        doc.add_heading(title, 0)

    for block in args.get("content", []):
        btype = block.get("type", "paragraph")
        if btype == "heading":
            doc.add_heading(block.get("text", ""), level=block.get("level", 1))
        elif btype == "paragraph":
            doc.add_paragraph(block.get("text", ""))
        elif btype == "bullet":
            for item in block.get("items", []):
                doc.add_paragraph(item, style="List Bullet")
        elif btype == "numbered":
            for item in block.get("items", []):
                doc.add_paragraph(item, style="List Number")
        elif btype == "table":
            headers = block.get("headers", [])
            rows = block.get("rows", [])
            if headers or rows:
                cols = max(len(headers), max((len(r) for r in rows), default=0))
                table = doc.add_table(rows=1 + len(rows), cols=cols)
                table.style = "Table Grid"
                if headers:
                    for j, h in enumerate(headers):
                        table.rows[0].cells[j].text = h
                for i, row in enumerate(rows):
                    for j, val in enumerate(row):
                        table.rows[i+1].cells[j].text = str(val)
        elif btype == "pagebreak":
            doc.add_page_break()

    doc.save(str(p))
    print(f"   Word document created: {p.name}")
    return {"status": "ok", "created": str(p), "type": "docx"}


def create_xlsx(args: Dict[str, Any], workspace: Path, safe_resolve) -> Dict[str, Any]:
    """
    Crea un file Excel (.xlsx).
    args: {
      "path": "desktop/foglio.xlsx",
      "sheets": [
        {
          "name": "Foglio1",
          "headers": ["Nome", "Età", "Città"],
          "rows": [["Mario", 30, "Roma"], ["Luca", 25, "Milano"]],
          "col_widths": [20, 10, 15]   (opzionale)
        }
      ]
    }
    """
    if not _check_lib("openpyxl"):
        return {"error": "missing_library", "reason": "Installa openpyxl: pip install openpyxl"}

    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    from openpyxl.utils import get_column_letter

    path = args.get("path", "")
    if not path:
        return {"error": "missing_path"}

    p = normalize_path(path, workspace)
    p.parent.mkdir(parents=True, exist_ok=True)

    wb = Workbook()
    wb.remove(wb.active)  # rimuove foglio default vuoto

    sheets = args.get("sheets", [])
    if not sheets:
        # fallback: un foglio vuoto
        sheets = [{"name": "Foglio1", "headers": [], "rows": []}]

    for sheet_def in sheets:
        ws = wb.create_sheet(title=sheet_def.get("name", "Foglio"))
        headers = sheet_def.get("headers", [])
        rows = sheet_def.get("rows", [])
        col_widths = sheet_def.get("col_widths", [])

        # Intestazioni con stile
        if headers:
            for j, h in enumerate(headers, 1):
                cell = ws.cell(row=1, column=j, value=h)
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill("solid", fgColor="2E75B6")
                cell.alignment = Alignment(horizontal="center")

        # Dati
        for i, row in enumerate(rows, 2 if headers else 1):
            for j, val in enumerate(row, 1):
                ws.cell(row=i, column=j, value=val)

        # Larghezze colonne
        for j, w in enumerate(col_widths, 1):
            ws.column_dimensions[get_column_letter(j)].width = w

        # Auto-width se non specificato
        if not col_widths:
            for col in ws.columns:
                max_len = max((len(str(cell.value or "")) for cell in col), default=8)
                ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 40)

    wb.save(str(p))
    print(f"   Excel file created: {p.name}")
    return {"status": "ok", "created": str(p), "type": "xlsx"}


def create_pptx(args: Dict[str, Any], workspace: Path, safe_resolve) -> Dict[str, Any]:
    """
    Crea un file PowerPoint (.pptx).
    args: {
      "path": "desktop/presentazione.pptx",
      "slides": [
        {"type": "title",   "title": "Titolo", "subtitle": "Sottotitolo"},
        {"type": "content", "title": "Slide 2", "bullets": ["Punto 1", "Punto 2"]},
        {"type": "blank",   "title": "Solo titolo"},
        {"type": "two_col", "title": "Due colonne", "left": ["A","B"], "right": ["C","D"]}
      ]
    }
    """
    if not _check_lib("pptx"):
        return {"error": "missing_library", "reason": "Installa python-pptx: pip install python-pptx"}

    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    path = args.get("path", "")
    if not path:
        return {"error": "missing_path"}

    p = normalize_path(path, workspace)
    p.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    LAYOUTS = {l.name: l for l in prs.slide_layouts}

    def _get_layout(name):
        return LAYOUTS.get(name) or prs.slide_layouts[1]

    for slide_def in args.get("slides", []):
        stype = slide_def.get("type", "content")

        if stype == "title":
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_def.get("title", "")
            if slide.placeholders[1]:
                slide.placeholders[1].text = slide_def.get("subtitle", "")

        elif stype == "content":
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_def.get("title", "")
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for bullet in slide_def.get("bullets", []):
                p_obj = tf.add_paragraph()
                p_obj.text = bullet
                p_obj.level = 0

        elif stype == "blank":
            layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(layout)
            if slide_def.get("title"):
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
                txBox.text_frame.text = slide_def.get("title", "")

        elif stype == "two_col":
            layout = prs.slide_layouts[3]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_def.get("title", "")
            left_items = slide_def.get("left", [])
            right_items = slide_def.get("right", [])
            placeholders = list(slide.placeholders)
            if len(placeholders) > 1:
                tf = placeholders[1].text_frame
                tf.clear()
                for item in left_items:
                    tf.add_paragraph().text = item
            if len(placeholders) > 2:
                tf = placeholders[2].text_frame
                tf.clear()
                for item in right_items:
                    tf.add_paragraph().text = item
        else:
            # fallback content
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_def.get("title", "Slide")

    prs.save(str(p))
    print(f"   PowerPoint created: {p.name}")
    return {"status": "ok", "created": str(p), "type": "pptx"}


def create_pdf(args: Dict[str, Any], workspace: Path, safe_resolve) -> Dict[str, Any]:
    """
    Crea un PDF.
    args: {
      "path": "desktop/documento.pdf",
      "title": "Titolo",
      "content": [
        {"type": "title",     "text": "Titolo grande"},
        {"type": "heading",   "text": "Sezione 1"},
        {"type": "paragraph", "text": "Testo normale..."},
        {"type": "bullet",    "items": ["voce 1", "voce 2"]},
        {"type": "table",     "headers": ["Col1","Col2"], "rows": [["A","B"]]},
        {"type": "spacer"}
      ]
    }
    """
    if not _check_lib("reportlab"):
        return {"error": "missing_library", "reason": "Installa reportlab: pip install reportlab"}

    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, ListFlowable, ListItem
    from reportlab.lib.enums import TA_LEFT, TA_CENTER

    path = args.get("path", "")
    if not path:
        return {"error": "missing_path"}

    p = normalize_path(path, workspace)
    p.parent.mkdir(parents=True, exist_ok=True)

    doc = SimpleDocTemplate(str(p), pagesize=A4,
                            leftMargin=2*cm, rightMargin=2*cm,
                            topMargin=2*cm, bottomMargin=2*cm)
    styles = getSampleStyleSheet()
    story = []

    title_style   = ParagraphStyle("PyTitle",   fontSize=22, spaceAfter=16, fontName="Helvetica-Bold", alignment=TA_CENTER)
    heading_style = ParagraphStyle("PyHeading", fontSize=14, spaceAfter=8,  fontName="Helvetica-Bold", textColor=colors.HexColor("#2E75B6"))
    body_style    = ParagraphStyle("PyBody",    fontSize=11, spaceAfter=6,  fontName="Helvetica",      leading=16)
    bullet_style  = ParagraphStyle("PyBullet",  fontSize=11, spaceAfter=4,  fontName="Helvetica",      leftIndent=20, bulletIndent=10)

    doc_title = args.get("title", "")
    if doc_title:
        story.append(Paragraph(doc_title, title_style))
        story.append(Spacer(1, 0.5*cm))

    for block in args.get("content", []):
        btype = block.get("type", "paragraph")
        if btype == "title":
            story.append(Paragraph(block.get("text", ""), title_style))
        elif btype == "heading":
            story.append(Paragraph(block.get("text", ""), heading_style))
        elif btype == "paragraph":
            story.append(Paragraph(block.get("text", ""), body_style))
        elif btype == "bullet":
            for item in block.get("items", []):
                story.append(Paragraph(f"• {item}", bullet_style))
        elif btype == "spacer":
            story.append(Spacer(1, 0.5*cm))
        elif btype == "table":
            headers = block.get("headers", [])
            rows = block.get("rows", [])
            data = ([headers] if headers else []) + rows
            if data:
                tbl = Table(data, hAlign="LEFT")
                tbl.setStyle(TableStyle([
                    ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#2E75B6")),
                    ("TEXTCOLOR",  (0,0), (-1,0), colors.white),
                    ("FONTNAME",   (0,0), (-1,0), "Helvetica-Bold"),
                    ("FONTSIZE",   (0,0), (-1,-1), 10),
                    ("GRID",       (0,0), (-1,-1), 0.5, colors.grey),
                    ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.white, colors.HexColor("#EBF3FB")]),
                    ("VALIGN",     (0,0), (-1,-1), "MIDDLE"),
                    ("PADDING",    (0,0), (-1,-1), 6),
                ]))
                story.append(tbl)
                story.append(Spacer(1, 0.3*cm))

    doc.build(story)
    print(f"   PDF created: {p.name}")
    return {"status": "ok", "created": str(p), "type": "pdf"}


# Aggiorna register_actions per includere i nuovi formati
_original_register = register_actions

def register_actions() -> Dict[str, Any]:
    actions = _original_register()
    actions.update({
        "create_docx": create_docx,
        "create_xlsx": create_xlsx,
        "create_pptx": create_pptx,
        "create_pdf":  create_pdf,
    })
    return actions